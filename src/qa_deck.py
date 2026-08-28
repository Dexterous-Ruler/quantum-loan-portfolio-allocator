"""Geometry QA for PRESENTATION.pptx.

This machine has no LibreOffice or poppler, so the deck cannot be rendered to images for
visual inspection. This checks the defects that rendering would have caught and that are
computable from the XML: shapes off the slide, shapes too close to the edge, overlapping
text boxes, and text that will not fit its container.

The text-fit estimate is deliberately conservative (assumes a wide average glyph), so it
over-reports rather than missing a real overflow.
"""
from __future__ import annotations

import math
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent.parent
DECK = ROOT / "PRESENTATION.pptx"

MARGIN_IN = 0.5          # brief's minimum breathing room from the slide edge
AVG_GLYPH_EM = 0.52      # average glyph width as a fraction of font size
LINE_SPACING = 1.22


def inches(v) -> float:
    return Emu(v).inches if v is not None else 0.0


def text_of(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs)


def max_font_pt(shape, default=12.0) -> float:
    sizes = []
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size is not None:
                sizes.append(r.font.size.pt)
    return max(sizes) if sizes else default


def estimate_height_in(text: str, width_in: float, font_pt: float) -> float:
    """Rough wrapped height. Counts explicit newlines as hard breaks."""
    if not text.strip() or width_in <= 0:
        return 0.0
    char_w_in = font_pt * AVG_GLYPH_EM / 72.0
    per_line = max(int(width_in / char_w_in), 1)
    lines = 0
    for para in text.split("\n"):
        lines += max(math.ceil(len(para) / per_line), 1)
    return lines * font_pt * LINE_SPACING / 72.0


def rects_overlap(a, b, tol=0.02) -> float:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ox = min(ax + aw, bx + bw) - max(ax, bx)
    oy = min(ay + ah, by + bh) - max(ay, by)
    if ox > tol and oy > tol:
        return ox * oy
    return 0.0


def main() -> int:
    # Slide text contains typographic minus/arrow glyphs that the Windows console's cp1252
    # codec cannot encode; without this the report crashes partway through.
    try:
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass

    prs = Presentation(DECK)
    SW, SH = prs.slide_width.inches, prs.slide_height.inches
    problems: list[str] = []

    for idx, slide in enumerate(prs.slides, 1):
        boxes = []
        for shape in slide.shapes:
            x, y = inches(shape.left), inches(shape.top)
            w, h = inches(shape.width), inches(shape.height)
            label = (text_of(shape).strip().replace("\n", " ")[:44] or shape.shape_type)

            if x < -0.01 or y < -0.01 or x + w > SW + 0.01 or y + h > SH + 0.01:
                problems.append(
                    f"slide {idx}: OFF-SLIDE  '{label}'  "
                    f"[{x:.2f},{y:.2f} {w:.2f}x{h:.2f}] vs {SW:.2f}x{SH:.2f}"
                )
            elif x < MARGIN_IN - 0.01 or y < MARGIN_IN - 0.01 \
                    or x + w > SW - MARGIN_IN + 0.01 or y + h > SH - MARGIN_IN + 0.01:
                problems.append(f"slide {idx}: TIGHT MARGIN  '{label}'  [{x:.2f},{y:.2f} {w:.2f}x{h:.2f}]")

            if shape.has_text_frame and text_of(shape).strip():
                need = estimate_height_in(text_of(shape), w, max_font_pt(shape))
                if need > h + 0.06:
                    problems.append(
                        f"slide {idx}: TEXT OVERFLOW  '{label}'  "
                        f"needs ~{need:.2f}in, box is {h:.2f}in"
                    )
                boxes.append(((x, y, w, h), label))

        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                area = rects_overlap(boxes[i][0], boxes[j][0])
                if area > 0.05:
                    problems.append(
                        f"slide {idx}: TEXT OVERLAP  '{boxes[i][1]}' x '{boxes[j][1]}'  "
                        f"({area:.2f} sq in)"
                    )

    print(f"{len(prs.slides)} slides, {SW:.2f}x{SH:.2f} in")
    if problems:
        print(f"\n{len(problems)} issue(s):")
        for p in problems:
            print("  " + p)
        return 1
    print("\nNo geometry issues found.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
